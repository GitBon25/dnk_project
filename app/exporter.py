from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.presentation_templates import PresentationTemplate, get_presentation_template
from app.schemas import PitchDeck


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def _set_background(slide, color: RGBColor) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text: str):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.clear()
    frame.paragraphs[0].text = text
    return box


def _style_textbox(
    box,
    *,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title_slide(prs: Presentation, deck: PitchDeck, theme: PresentationTemplate) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, _rgb(theme.background))

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.7),
        Inches(0.75),
        Inches(0.18),
        Inches(5.9),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _rgb(theme.accent)
    accent_bar.line.fill.background()

    eyebrow = _add_textbox(
        slide,
        Inches(1.15),
        Inches(0.82),
        Inches(4.4),
        Inches(0.4),
        deck.audience.upper(),
    )
    _style_textbox(
        eyebrow,
        font_name=theme.body_font,
        font_size=12,
        color=_rgb(theme.accent_dark),
        bold=True,
    )

    title = _add_textbox(
        slide,
        Inches(1.15),
        Inches(1.35),
        Inches(7.8),
        Inches(1.5),
        deck.startup_name,
    )
    _style_textbox(
        title,
        font_name=theme.title_font,
        font_size=28,
        color=_rgb(theme.ink),
        bold=True,
    )

    subtitle = _add_textbox(
        slide,
        Inches(1.15),
        Inches(2.45),
        Inches(6.7),
        Inches(1.6),
        deck.one_liner,
    )
    _style_textbox(
        subtitle,
        font_name=theme.body_font,
        font_size=18,
        color=_rgb(theme.muted),
    )

    info_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.95),
        Inches(0.9),
        Inches(3.45),
        Inches(5.75),
    )
    info_card.fill.solid()
    info_card.fill.fore_color.rgb = _rgb(theme.surface)
    info_card.line.color.rgb = _rgb(theme.line)

    info_title = _add_textbox(
        slide,
        Inches(9.25),
        Inches(1.2),
        Inches(2.8),
        Inches(0.5),
        "Deck Snapshot",
    )
    _style_textbox(
        info_title,
        font_name=theme.title_font,
        font_size=18,
        color=_rgb(theme.ink),
        bold=True,
    )

    gaps_label = _add_textbox(
        slide,
        Inches(9.25),
        Inches(1.85),
        Inches(2.7),
        Inches(0.35),
        "Open Questions",
    )
    _style_textbox(
        gaps_label,
        font_name=theme.body_font,
        font_size=10,
        color=_rgb(theme.muted),
        bold=True,
    )

    gaps_frame = slide.shapes.add_textbox(Inches(9.2), Inches(2.15), Inches(2.85), Inches(3.85)).text_frame
    gaps_frame.word_wrap = True
    gaps_frame.clear()

    if deck.info_gaps:
        for index, gap in enumerate(deck.info_gaps[:4]):
            paragraph = gaps_frame.paragraphs[0] if index == 0 else gaps_frame.add_paragraph()
            paragraph.text = gap
            paragraph.level = 0
            paragraph.space_after = Pt(10)
            paragraph.bullet = True
            run = paragraph.runs[0]
            run.font.name = theme.body_font
            run.font.size = Pt(11)
            run.font.color.rgb = _rgb(theme.ink)
    else:
        gaps_frame.paragraphs[0].text = "No major information gaps surfaced."
        run = gaps_frame.paragraphs[0].runs[0]
        run.font.name = theme.body_font
        run.font.size = Pt(11)
        run.font.color.rgb = _rgb(theme.muted)


def _add_content_slide(
    prs: Presentation,
    slide_number: int,
    title_text: str,
    bullets: list[str],
    theme: PresentationTemplate,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, _rgb(theme.background))

    top_rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.65),
        Inches(0.55),
        Inches(12.0),
        Inches(0.08),
    )
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = _rgb(theme.accent)
    top_rule.line.fill.background()

    slide_num = _add_textbox(
        slide,
        Inches(0.7),
        Inches(0.78),
        Inches(0.6),
        Inches(0.4),
        f"{slide_number:02d}",
    )
    _style_textbox(
        slide_num,
        font_name=theme.body_font,
        font_size=11,
        color=_rgb(theme.accent_dark),
        bold=True,
    )

    title = _add_textbox(
        slide,
        Inches(0.7),
        Inches(1.15),
        Inches(5.8),
        Inches(0.8),
        title_text,
    )
    _style_textbox(
        title,
        font_name=theme.title_font,
        font_size=24,
        color=_rgb(theme.ink),
        bold=True,
    )

    subtitle = _add_textbox(
        slide,
        Inches(0.7),
        Inches(1.9),
        Inches(5.0),
        Inches(0.45),
        "Core talking points for the slide",
    )
    _style_textbox(
        subtitle,
        font_name=theme.body_font,
        font_size=11,
        color=_rgb(theme.muted),
    )

    positions = [
        (Inches(0.75), Inches(2.55)),
        (Inches(4.45), Inches(2.55)),
        (Inches(8.15), Inches(2.55)),
    ]

    for index, bullet in enumerate(bullets[:3]):
        left, top = positions[index]
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            Inches(3.0),
            Inches(3.35),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(theme.surface)
        card.line.color.rgb = _rgb(theme.line)

        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            left + Inches(0.22),
            top + Inches(0.22),
            Inches(0.42),
            Inches(0.42),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb(theme.accent)
        badge.line.fill.background()

        badge_text = _add_textbox(
            slide,
            left + Inches(0.22),
            top + Inches(0.215),
            Inches(0.42),
            Inches(0.28),
            str(index + 1),
        )
        _style_textbox(
            badge_text,
            font_name=theme.body_font,
            font_size=10,
            color=RGBColor(255, 255, 255),
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        bullet_box = _add_textbox(
            slide,
            left + Inches(0.3),
            top + Inches(0.88),
            Inches(2.4),
            Inches(2.1),
            bullet,
        )
        _style_textbox(
            bullet_box,
            font_name=theme.body_font,
            font_size=17,
            color=_rgb(theme.ink),
            bold=True,
        )


def create_ppt(deck: PitchDeck, output_path: str, template: str | None = None) -> str:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    theme = get_presentation_template(template)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _add_title_slide(prs, deck, theme)

    for slide_number, slide_data in enumerate(deck.slides[1:], start=2):
        _add_content_slide(prs, slide_number, slide_data.title, slide_data.bullets, theme)

    prs.save(output)
    return str(output)
